"""Convert PPTX to PDF by rendering each slide using python-pptx + reportlab + Pillow.

This script extracts shapes from each slide and renders them onto PDF pages
with correct positioning, colors, fonts, and text formatting.
"""

import sys
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from reportlab.lib.pagesizes import landscape
from reportlab.lib.units import inch
from reportlab.pdfgen import canvas
from reportlab.lib.colors import Color, HexColor, white, black
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.ttfonts import TTFont
from io import BytesIO
from PIL import Image as PILImage
import tempfile


def emu_to_points(emu):
    """Convert EMU to points (1 inch = 914400 EMU = 72 points)."""
    return emu * 72.0 / 914400.0


def rgb_to_color(rgb):
    """Convert python-pptx RGBColor to reportlab Color."""
    if rgb is None:
        return None
    return Color(rgb[0] / 255.0, rgb[1] / 255.0, rgb[2] / 255.0)


def get_shape_fill_color(shape):
    """Extract fill color from a shape."""
    try:
        fill = shape.fill
        if fill.type is not None:
            try:
                fc = fill.fore_color
                if fc and fc.type is not None:
                    return rgb_to_color(fc.rgb)
            except:
                pass
    except:
        pass
    return None


def get_text_color(run):
    """Extract text color from a run."""
    try:
        if run.font.color and run.font.color.type is not None:
            return rgb_to_color(run.font.color.rgb)
    except:
        pass
    return black


def draw_shape(c, shape, slide_w, slide_h, page_w, page_h):
    """Draw a shape on the canvas."""
    if not hasattr(shape, 'left') or shape.left is None:
        return

    # Scale factors
    sx = page_w / slide_w
    sy = page_h / slide_h

    x = emu_to_points(shape.left) * sx
    y_top = emu_to_points(shape.top) * sy
    w = emu_to_points(shape.width) * sx
    h = emu_to_points(shape.height) * sy

    # Convert to reportlab coords (origin at bottom-left)
    y = page_h - y_top - h

    # Draw fill
    fill_color = get_shape_fill_color(shape)
    if fill_color:
        c.setFillColor(fill_color)
        c.rect(x, y, w, h, fill=1, stroke=0)

    # Draw border
    try:
        line = shape.line
        if line and line.fill and line.fill.type is not None:
            try:
                lc = line.color.rgb
                c.setStrokeColor(rgb_to_color(lc))
                lw = emu_to_points(line.width) if line.width else 1
                c.setLineWidth(lw * sx)
                c.rect(x, y, w, h, fill=0, stroke=1)
            except:
                pass
    except:
        pass

    # Draw image
    if shape.shape_type == 13:  # Picture
        try:
            img_blob = shape.image.blob
            img_stream = BytesIO(img_blob)
            pil_img = PILImage.open(img_stream)

            # Save to temp file
            with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp:
                if pil_img.mode != 'RGBA':
                    pil_img = pil_img.convert('RGBA')
                pil_img.save(tmp, 'PNG')
                tmp_path = tmp.name

            from reportlab.lib.utils import ImageReader
            c.drawImage(tmp_path, x, y, w, h, preserveAspectRatio=False, mask='auto')
            os.unlink(tmp_path)
        except Exception as e:
            pass

    # Draw text
    if shape.has_text_frame:
        tf = shape.text_frame
        # Padding
        margin_l = emu_to_points(tf.margin_left or 0) * sx
        margin_r = emu_to_points(tf.margin_right or 0) * sx
        margin_t = emu_to_points(tf.margin_top or 0) * sy
        margin_b = emu_to_points(tf.margin_bottom or 0) * sy

        text_x = x + margin_l
        text_w = w - margin_l - margin_r
        text_y_start = page_h - y_top - margin_t

        cursor_y = text_y_start

        for para in tf.paragraphs:
            if not para.text.strip():
                cursor_y -= 10 * sy  # blank line
                continue

            # Determine alignment
            align = para.alignment

            for run in para.runs:
                if not run.text:
                    continue

                # Font settings
                font_size = 10
                if run.font.size:
                    font_size = run.font.size.pt * sy
                font_size = max(font_size, 4)

                is_bold = run.font.bold
                font_name = 'Helvetica'
                if is_bold:
                    font_name = 'Helvetica-Bold'

                text_color = get_text_color(run)
                c.setFillColor(text_color)
                c.setFont(font_name, font_size)

                # Position text
                cursor_y -= font_size * 1.15

                text = run.text

                # Simple word wrap
                words = text.split(' ')
                line = ''
                for word in words:
                    test_line = line + (' ' if line else '') + word
                    tw = c.stringWidth(test_line, font_name, font_size)
                    if tw > text_w and line:
                        # Draw current line
                        if align == PP_ALIGN.CENTER:
                            lw = c.stringWidth(line, font_name, font_size)
                            c.drawString(text_x + (text_w - lw) / 2, cursor_y, line)
                        elif align == PP_ALIGN.RIGHT:
                            lw = c.stringWidth(line, font_name, font_size)
                            c.drawString(text_x + text_w - lw, cursor_y, line)
                        else:
                            c.drawString(text_x, cursor_y, line)
                        cursor_y -= font_size * 1.2
                        line = word
                    else:
                        line = test_line

                if line:
                    if align == PP_ALIGN.CENTER:
                        lw = c.stringWidth(line, font_name, font_size)
                        c.drawString(text_x + (text_w - lw) / 2, cursor_y, line)
                    elif align == PP_ALIGN.RIGHT:
                        lw = c.stringWidth(line, font_name, font_size)
                        c.drawString(text_x + text_w - lw, cursor_y, line)
                    else:
                        c.drawString(text_x, cursor_y, line)

            cursor_y -= 2 * sy  # paragraph spacing

    # Draw table
    if shape.has_table:
        table = shape.table
        rows = len(table.rows)
        cols = len(table.columns)

        # Calculate column widths
        col_widths = []
        for col in table.columns:
            col_widths.append(emu_to_points(col.width) * sx)

        row_heights = []
        for row in table.rows:
            row_heights.append(emu_to_points(row.height) * sy)

        cell_y = page_h - y_top
        for ri, row in enumerate(table.rows):
            cell_x = x
            rh = row_heights[ri]
            cell_y -= rh

            for ci, cell in enumerate(row.cells):
                cw = col_widths[ci]

                # Cell fill
                try:
                    cf = cell.fill
                    if cf.type is not None:
                        cc = rgb_to_color(cf.fore_color.rgb)
                        if cc:
                            c.setFillColor(cc)
                            c.rect(cell_x, cell_y, cw, rh, fill=1, stroke=0)
                except:
                    pass

                # Cell border
                c.setStrokeColor(Color(0.8, 0.8, 0.8))
                c.setLineWidth(0.5)
                c.rect(cell_x, cell_y, cw, rh, fill=0, stroke=1)

                # Cell text
                text_y = cell_y + rh - 4 * sy
                for para in cell.text_frame.paragraphs:
                    for run in para.runs:
                        if not run.text:
                            continue
                        fs = 7 * sy
                        if run.font.size:
                            fs = run.font.size.pt * sy
                        fs = max(fs, 4)

                        fn = 'Helvetica-Bold' if run.font.bold else 'Helvetica'
                        tc = get_text_color(run)

                        c.setFillColor(tc)
                        c.setFont(fn, fs)
                        text_y -= fs * 1.1

                        # Truncate if too wide
                        text = run.text
                        while c.stringWidth(text, fn, fs) > cw - 4 * sx and len(text) > 1:
                            text = text[:-1]

                        c.drawString(cell_x + 2 * sx, text_y, text)

                cell_x += cw


def convert_pptx_to_pdf(pptx_path, pdf_path):
    """Convert a PPTX file to PDF."""
    prs = Presentation(pptx_path)

    slide_w = prs.slide_width
    slide_h = prs.slide_height

    # Page size in points (16:9 widescreen)
    page_w = 720
    page_h = page_w * (slide_h / slide_w)

    c = canvas.Canvas(pdf_path, pagesize=(page_w, page_h))

    for i, slide in enumerate(prs.slides):
        print(f"  Rendering slide {i + 1}...")

        # Slide background
        bg = slide.background
        bg_color = None
        try:
            if bg.fill.type is not None:
                bg_color = rgb_to_color(bg.fill.fore_color.rgb)
        except:
            pass

        if bg_color:
            c.setFillColor(bg_color)
        else:
            c.setFillColor(white)
        c.rect(0, 0, page_w, page_h, fill=1, stroke=0)

        # Draw shapes in order
        for shape in slide.shapes:
            try:
                draw_shape(c, shape, slide_w, slide_h, page_w, page_h)
            except Exception as e:
                pass

        c.showPage()

    c.save()
    print(f"  PDF saved: {pdf_path}")


if __name__ == '__main__':
    pptx_path = sys.argv[1] if len(sys.argv) > 1 else '/Users/ilkerkadirozturk/Documents/genomer_brochures/sunum/Genomer_Microbiome_Testing_Line.pptx'
    pdf_path = sys.argv[2] if len(sys.argv) > 2 else pptx_path.replace('.pptx', '.pdf')

    print(f"Converting: {pptx_path}")
    convert_pptx_to_pdf(pptx_path, pdf_path)
